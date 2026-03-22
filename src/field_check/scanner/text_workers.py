"""Worker functions for text extraction (ProcessPoolExecutor).

These must be top-level functions in a module so they can be pickled
by ProcessPoolExecutor. Separated from text.py to keep file sizes
manageable.
"""

from __future__ import annotations

import os

from field_check.scanner.text import (
    _MAX_TEXT_READ,
    CHARS_PER_PAGE_IMAGE_HEAVY,
    CHARS_PER_PAGE_TEXT_HEAVY,
    CLASSIFICATION_IMAGE_HEAVY,
    CLASSIFICATION_MIXED,
    CLASSIFICATION_TEXT_HEAVY,
    TEXT_SIZE_RATIO_IMAGE_HEAVY,
    TextResult,
)


def _docx_full_text(doc: object) -> str:
    """Extract all text from a python-docx Document.

    Includes body paragraphs, table cells, headers, and footers.
    """
    parts: list[str] = []

    # Body paragraphs
    for p in doc.paragraphs:  # type: ignore[attr-defined]
        if p.text:
            parts.append(p.text)

    # Table cells (deduplicate merged cells by XML element identity)
    seen_cells: set[int] = set()
    for table in doc.tables:  # type: ignore[attr-defined]
        for row in table.rows:
            for cell in row.cells:
                cell_id = id(cell._tc)  # type: ignore[attr-defined]
                if cell_id not in seen_cells and cell.text:
                    seen_cells.add(cell_id)
                    parts.append(cell.text)

    # Headers and footers
    for section in doc.sections:  # type: ignore[attr-defined]
        for hf in (section.header, section.footer):
            if hf and hf.is_linked_to_previous is False:
                for p in hf.paragraphs:
                    if p.text:
                        parts.append(p.text)

    return "\n".join(parts)


def _extract_pdf(filepath: str) -> TextResult:
    """Extract text, metadata, and classification from a PDF.

    Single-pass extraction: opens the file once to get text, page count,
    scanned detection (via page.chars), content classification, and metadata.
    """
    import pdfplumber

    result = TextResult(path=filepath)
    file_size = os.path.getsize(filepath)

    try:
        with pdfplumber.open(filepath) as pdf:
            result.page_count = len(pdf.pages)

            # Extract metadata
            meta = pdf.metadata or {}
            result.metadata["title"] = meta.get("Title") or None
            result.metadata["author"] = meta.get("Author") or None
            creation = meta.get("CreationDate")
            result.metadata["creation_date"] = str(creation) if creation else None

            total_chars = 0
            total_text_bytes = 0
            scanned_pages = 0
            native_pages = 0

            page_texts: list[str] = []
            high_image_pages = 0
            for page in pdf.pages:
                # Count char objects for scanned detection
                char_count = len(page.chars) if page.chars else 0
                if char_count == 0:
                    scanned_pages += 1
                else:
                    native_pages += 1
                    total_chars += char_count

                # Image coverage heuristic: check if page is image-dominated
                try:
                    images = page.images if hasattr(page, "images") else []
                    page_w = float(page.width) if page.width else 1.0
                    page_h = float(page.height) if page.height else 1.0
                    page_area = page_w * page_h
                    img_area = sum(
                        float(img.get("width", 0)) * float(img.get("height", 0)) for img in images
                    )
                    if page_area > 0 and img_area / page_area >= 0.9:
                        high_image_pages += 1
                except Exception:
                    pass  # Image analysis is best-effort

                # Extract text
                page_text = page.extract_text() or ""
                page_texts.append(page_text)
                total_text_bytes += len(page_text.encode("utf-8", errors="replace"))

            result.text = "\n".join(page_texts)
            result.text_length = len(result.text)

            # Scanned detection — enhanced with image coverage heuristic
            if result.page_count > 0:
                chars_pp = total_chars / result.page_count
                # OCR'd scanned PDFs: have text but are image-dominated
                if (
                    high_image_pages >= result.page_count * 0.8 and chars_pp < 200
                ) or scanned_pages == result.page_count:
                    result.is_scanned = True
                elif scanned_pages > 0 and native_pages > 0:
                    result.is_mixed_scan = True

            # Text quality check: high U+FFFD ratio indicates extraction issues
            if result.text_length > 0:
                fffd_count = result.text.count("\ufffd")
                if fffd_count > result.text_length * 0.10:
                    result.extraction_quality = "degraded"
                else:
                    result.extraction_quality = "good"

            # Content classification
            if result.page_count > 0:
                result.chars_per_page = total_chars / result.page_count
                result.text_size_ratio = total_text_bytes / file_size if file_size > 0 else 0.0

                if result.is_scanned or result.chars_per_page < CHARS_PER_PAGE_IMAGE_HEAVY:
                    result.classification = CLASSIFICATION_IMAGE_HEAVY
                elif result.chars_per_page > CHARS_PER_PAGE_TEXT_HEAVY:
                    result.classification = CLASSIFICATION_TEXT_HEAVY
                else:
                    # Mixed zone -- check secondary metric
                    if result.text_size_ratio < TEXT_SIZE_RATIO_IMAGE_HEAVY:
                        result.classification = CLASSIFICATION_IMAGE_HEAVY
                    else:
                        result.classification = CLASSIFICATION_MIXED

    except Exception as exc:
        result.error = str(exc)

    return result


def _extract_docx(filepath: str) -> TextResult:
    """Extract text and metadata from a DOCX file."""
    from docx import Document

    result = TextResult(path=filepath)
    file_size = os.path.getsize(filepath)

    try:
        doc = Document(filepath)

        # Extract text from paragraphs, tables, headers, footers
        text = _docx_full_text(doc)
        result.text = text
        text_bytes = len(text.encode("utf-8", errors="replace"))
        result.text_length = len(text)
        result.text_size_ratio = text_bytes / file_size if file_size > 0 else 0.0

        # Extract metadata
        props = doc.core_properties
        result.metadata["title"] = props.title if props.title else None
        result.metadata["author"] = props.author if props.author else None
        result.metadata["creation_date"] = props.created.isoformat() if props.created else None

        # DOCX is always text-based, never scanned
        result.classification = CLASSIFICATION_TEXT_HEAVY

    except Exception as exc:
        result.error = str(exc)

    return result


def _extract_xlsx(filepath: str) -> TextResult:
    """Extract text from an XLSX file using openpyxl (if installed)."""
    result = TextResult(path=filepath)
    try:
        import openpyxl

        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        try:
            parts: list[str] = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        parts.append("\t".join(cells))
            result.text = "\n".join(parts)
            result.text_length = len(result.text)
            result.classification = CLASSIFICATION_TEXT_HEAVY
        finally:
            wb.close()
    except ImportError:
        result.error = "openpyxl not installed (pip install field-check[formats])"
    except Exception as exc:
        result.error = str(exc)
    return result


def _extract_pptx(filepath: str) -> TextResult:
    """Extract text from a PPTX file using python-pptx (if installed)."""
    result = TextResult(path=filepath)
    try:
        from pptx import Presentation

        prs = Presentation(filepath)
        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text:
                            parts.append(para.text)
        result.text = "\n".join(parts)
        result.text_length = len(result.text)
        result.page_count = len(prs.slides)
        result.classification = CLASSIFICATION_TEXT_HEAVY
    except ImportError:
        result.error = "python-pptx not installed (pip install field-check[formats])"
    except Exception as exc:
        result.error = str(exc)
    return result


def _extract_eml(filepath: str) -> TextResult:
    """Extract text from an EML (RFC 822) email file using stdlib."""
    import email
    import email.policy

    result = TextResult(path=filepath)
    try:
        with open(filepath, "rb") as f:
            msg = email.message_from_binary_file(f, policy=email.policy.default)
        parts: list[str] = []
        for header in ("Subject", "From", "To", "Date"):
            val = msg.get(header)
            if val:
                parts.append(f"{header}: {val}")
        body = msg.get_body(preferencelist=("plain", "html"))
        if body:
            content = body.get_content()
            if isinstance(content, str):
                parts.append(content)
        result.text = "\n".join(parts)
        result.text_length = len(result.text)
        result.metadata["title"] = msg.get("Subject") or None
        result.metadata["author"] = msg.get("From") or None
        result.metadata["creation_date"] = msg.get("Date") or None
        result.classification = CLASSIFICATION_TEXT_HEAVY
    except Exception as exc:
        result.error = str(exc)
    return result


def _extract_epub(filepath: str) -> TextResult:
    """Extract text from an EPUB file using stdlib zipfile + html.parser."""
    import re
    import zipfile
    from html.parser import HTMLParser

    class _TextExtractor(HTMLParser):
        def __init__(self) -> None:
            super().__init__()
            self.parts: list[str] = []

        def handle_data(self, data: str) -> None:
            stripped = data.strip()
            if stripped:
                self.parts.append(stripped)

    result = TextResult(path=filepath)
    try:
        with zipfile.ZipFile(filepath) as zf:
            parts: list[str] = []
            total_bytes = 0
            for name in zf.namelist():
                if re.search(r"\.(xhtml|html|htm)$", name, re.IGNORECASE):
                    raw_bytes = zf.read(name)
                    total_bytes += len(raw_bytes)
                    if total_bytes > _MAX_TEXT_READ:
                        break  # size limit to prevent zip bombs
                    raw = raw_bytes.decode("utf-8", errors="replace")
                    parser = _TextExtractor()
                    parser.feed(raw)
                    parts.extend(parser.parts)
        result.text = "\n".join(parts)
        result.text_length = len(result.text)
        result.classification = CLASSIFICATION_TEXT_HEAVY
    except Exception as exc:
        result.error = str(exc)
    return result


_XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _extract_single(filepath: str, mime_type: str) -> TextResult:
    """Worker function for ProcessPoolExecutor.

    Must be top-level for pickling. Dispatches to type-specific extractors.
    """
    if mime_type == "application/pdf":
        return _extract_pdf(filepath)
    if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return _extract_docx(filepath)
    if mime_type == _XLSX_MIME:
        return _extract_xlsx(filepath)
    if mime_type == _PPTX_MIME:
        return _extract_pptx(filepath)
    if mime_type == "message/rfc822":
        return _extract_eml(filepath)
    if mime_type == "application/epub+zip":
        return _extract_epub(filepath)
    return TextResult(path=filepath, error=f"Unsupported type: {mime_type}")


def _extract_pdf_text_fast(filepath: str) -> tuple[str, str | None, float, str | None]:
    """Extract PDF text using pdf_oxide (fast) with pdfplumber fallback."""
    try:
        from pdf_oxide import PdfDocument

        doc = PdfDocument(filepath)
        text = doc.to_plain_text_all()[:_MAX_TEXT_READ]
        if text.strip():
            return (text, None, 0.0, None)
        # Empty text — fall through to pdfplumber for better extraction
    except ImportError:
        pass
    except Exception:
        pass  # Fall through to pdfplumber

    import pdfplumber

    with pdfplumber.open(filepath) as pdf:
        text = "\n".join(page.extract_text() or "" for page in pdf.pages)
    return (text, None, 0.0, None)


def _extract_text_for_cache(
    filepath: str, mime_type: str
) -> tuple[str, str | None, float, str | None]:
    """Worker function to extract text content from any supported file.

    Must be top-level for pickling (ProcessPoolExecutor).

    Returns:
        Tuple of (text, encoding_name, encoding_confidence, error).
        encoding_name is only set for plain text files.
    """
    try:
        if mime_type == "application/pdf":
            return _extract_pdf_text_fast(filepath)

        if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            from docx import Document

            doc = Document(filepath)
            text = _docx_full_text(doc)
            return (text, None, 0.0, None)

        if mime_type == _XLSX_MIME:
            r = _extract_xlsx(filepath)
            return (r.text, None, 0.0, r.error)

        if mime_type == _PPTX_MIME:
            r = _extract_pptx(filepath)
            return (r.text, None, 0.0, r.error)

        if mime_type == "message/rfc822":
            r = _extract_eml(filepath)
            return (r.text, None, 0.0, r.error)

        if mime_type == "application/epub+zip":
            r = _extract_epub(filepath)
            return (r.text, None, 0.0, r.error)

        # Plain text types — read bytes, detect encoding
        with open(filepath, "rb") as f:
            raw = f.read(_MAX_TEXT_READ)

        text, enc, conf = _detect_encoding(raw)
        return (text, enc, conf, None)
    except Exception as exc:
        return ("", None, 0.0, str(exc))


def _detect_encoding(raw: bytes) -> tuple[str, str, float]:
    """Detect encoding of raw bytes using charset-normalizer.

    Returns:
        Tuple of (decoded_text, encoding_name, confidence).
    """
    try:
        from charset_normalizer import from_bytes

        result = from_bytes(raw).best()
        if result:
            return (str(result), result.encoding, 1.0 - result.chaos)
    except ImportError:
        pass
    except Exception:
        pass

    # Fallback to chardet if available (optional)
    try:
        import chardet

        detection = chardet.detect(raw)
        enc = detection.get("encoding")
        conf = detection.get("confidence", 0.0) or 0.0
        if enc and conf > 0.5:
            text = raw.decode(enc, errors="replace")
            return (text, enc, conf)
    except ImportError:
        pass
    except Exception:
        pass

    # Last resort: assume UTF-8
    return (raw.decode("utf-8", errors="replace"), "utf-8", 0.0)


def _extract_plain_text(
    filepath: str,
) -> tuple[str, str | None, float, str | None]:
    """Extract text and encoding from a plain text file.

    Must be top-level for pickling (ProcessPoolExecutor).

    Returns:
        Tuple of (text, encoding_name, encoding_confidence, error).
    """
    try:
        with open(filepath, "rb") as f:
            raw = f.read(_MAX_TEXT_READ)

        text, enc, conf = _detect_encoding(raw)
        return (text, enc, conf, None)
    except Exception as exc:
        return ("", None, 0.0, str(exc))
